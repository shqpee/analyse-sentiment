# -*- coding: utf-8 -*-
"""
generer_presentation.py — Génère un PowerPoint structuré (RNCP37827) pour le
projet « Analyseur de sentiment d'avis clients », inspiré du deck RouteZone.

Sortie : Presentation_Avis_Sentiment_AMELIOREE.pptx
Lancer : python generer_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Thème / couleurs ─────────────────────────────────────────────
NAVY     = RGBColor(0x1E, 0x1B, 0x4B)   # fond des slides de section
INDIGO   = RGBColor(0x63, 0x66, 0xF1)
VIOLET   = RGBColor(0x8B, 0x5C, 0xF6)
TEAL     = RGBColor(0x14, 0xB8, 0xA6)
ORANGE   = RGBColor(0xF5, 0x9E, 0x0B)
CRIMSON  = RGBColor(0xE1, 0x1D, 0x48)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
TEXTDARK = RGBColor(0x1F, 0x29, 0x37)
GREY     = RGBColor(0x6B, 0x72, 0x80)
LIGHT    = RGBColor(0xEE, 0xF0, 0xFB)   # fond des cartes claires
CODEBG   = RGBColor(0x1E, 0x1B, 0x4B)
CODEGREEN= RGBColor(0x86, 0xEF, 0xAC)
PLACE    = RGBColor(0xE5, 0xE7, 0xEB)   # fond encart image
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# Couleur par épreuve (comme RouteZone)
EP_COLOR = {"E1": INDIGO, "E2": ORANGE, "E3": VIOLET, "E4": TEAL, "E5": CRIMSON}

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


# ── Helpers ──────────────────────────────────────────────────────
def _set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    return p


def rect(slide, x, y, w, h, fill, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for (text, size, color, bold, align, *rest) in lines:
        font = rest[0] if rest else "Calibri"
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def eyebrow_title(slide, eyebrow, title, color):
    """Bandeau du haut : petit eyebrow coloré + grand titre, comme RouteZone."""
    textbox(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4),
            [(eyebrow, 13, color, True, PP_ALIGN.LEFT)])
    textbox(slide, Inches(0.6), Inches(0.7), Inches(12.1), Inches(0.9),
            [(title, 30, TEXTDARK, True, PP_ALIGN.LEFT)])


def chip(slide, label, color, x=Inches(11.7), y=Inches(0.35)):
    sp = rect(slide, x, y, Inches(1.0), Inches(0.5), color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, label, 14, WHITE, bold=True, align=PP_ALIGN.CENTER)


def card(slide, x, y, w, h, title, bullets, accent=INDIGO, title_color=None,
         body_color=TEXTDARK, fill=LIGHT, badge=None):
    """Carte claire avec titre + puces (style RouteZone)."""
    box = rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18),
                                  w - Inches(0.44), h - Inches(0.36))
    tf = tb.text_frame
    tf.word_wrap = True
    # Titre
    p = tf.paragraphs[0]
    if badge:
        rb = p.add_run(); rb.text = badge + "  "
        rb.font.size = Pt(13); rb.font.bold = True; rb.font.color.rgb = accent
    rt = p.add_run(); rt.text = title
    rt.font.size = Pt(15); rt.font.bold = True
    rt.font.color.rgb = title_color or accent
    # Puces
    for b in bullets:
        pp = tf.add_paragraph()
        pp.space_before = Pt(5)
        r = pp.add_run(); r.text = "• " + b
        r.font.size = Pt(12); r.font.color.rgb = body_color
    return box


def code_box(slide, x, y, w, h, title, code):
    box = rect(slide, x, y, w, h, CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.2),
                                  w - Inches(0.5), h - Inches(0.4))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = WHITE
    for line in code.split("\n"):
        pp = tf.add_paragraph()
        r = pp.add_run(); r.text = line if line else " "
        r.font.size = Pt(11); r.font.name = "Consolas"; r.font.color.rgb = CODEGREEN
    return box


def image_placeholder(slide, x, y, w, h, label):
    """Encart vide : zone grise pour insérer une image plus tard."""
    box = rect(slide, x, y, w, h, PLACE, line=GREY, line_w=Pt(1.5),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "🖼  ENCART IMAGE"
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = GREY
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = label
    r2.font.size = Pt(11); r2.font.italic = True; r2.font.color.rgb = GREY
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = "(à insérer)"
    r3.font.size = Pt(9); r3.font.color.rgb = GREY
    return box


def def_box(slide, x, y, w, h, text):
    """Encadré 'définition' bleu clair (📖) comme RouteZone."""
    box = rect(slide, x, y, w, h, RGBColor(0xE8, 0xEE, 0xFB),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1),
                                  w - Inches(0.4), h - Inches(0.2))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "📖 " + text
    r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = RGBColor(0x37,0x41,0x51)
    return box


def banner(slide, text, color, y=Inches(6.7)):
    box = rect(slide, Inches(0.6), y, Inches(12.13), Inches(0.55), color,
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = box.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, text, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)


def footer(slide, n):
    textbox(slide, Inches(12.2), Inches(7.05), Inches(1.0), Inches(0.3),
            [(str(n), 10, GREY, False, PP_ALIGN.RIGHT)])


_page = {"n": 0}
def new_slide():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, WHITE)  # fond blanc
    _page["n"] += 1
    return s


def section_slide(ep, title, sub):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, NAVY)
    _page["n"] += 1
    # gros bloc E?
    blk = rect(s, Inches(0.9), Inches(2.6), Inches(2.0), Inches(1.3),
               EP_COLOR[ep], shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = blk.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, ep, 44, WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, Inches(3.3), Inches(2.7), Inches(9), Inches(1.2),
            [(title, 30, WHITE, True, PP_ALIGN.LEFT)])
    textbox(s, Inches(3.32), Inches(3.9), Inches(9), Inches(0.6),
            [(sub, 15, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.LEFT)])
    footer(s, _page["n"])
    return s


# =================================================================
# 1. TITRE
# =================================================================
s = prs.slides.add_slide(BLANK); _page["n"] += 1
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 0, Inches(0.25), SH, INDIGO)
textbox(s, Inches(0.9), Inches(0.7), Inches(11), Inches(0.5),
        [("TITRE PROFESSIONNEL RNCP37827  ·  NIVEAU 6 (BAC+3/4)", 13, INDIGO, True, PP_ALIGN.LEFT)])
textbox(s, Inches(0.9), Inches(1.1), Inches(11), Inches(0.5),
        [("Développeur en Intelligence Artificielle", 17, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.LEFT)])
textbox(s, Inches(0.85), Inches(2.6), Inches(11.5), Inches(1.4),
        [("Analyseur de sentiment d'avis clients", 42, WHITE, True, PP_ALIGN.LEFT)])
rect(s, Inches(0.95), Inches(3.9), Inches(2.5), Inches(0.06), INDIGO)
textbox(s, Inches(0.95), Inches(4.6), Inches(6), Inches(0.5),
        [("Caroline Bacheron", 20, WHITE, True, PP_ALIGN.LEFT)])
textbox(s, Inches(0.95), Inches(5.2), Inches(7), Inches(1.2),
        [("Pipeline complet : collecte → modèle IA → application web",
          14, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.LEFT),
         ("Modèle de production : XGBoost (TF-IDF) entraîné en français",
          14, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.LEFT)])
image_placeholder(s, Inches(9.7), Inches(2.5), Inches(2.7), Inches(2.7), "Photo / logo Simplon")
textbox(s, Inches(9.7), Inches(5.3), Inches(3), Inches(1),
        [("Accompagnée par :", 12, INDIGO, True, PP_ALIGN.LEFT),
         ("W. Ben Chelbi · A. Sakhri · G. Soulat", 12, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.LEFT)])

# =================================================================
# 2. LE PROJET EN UNE PHRASE
# =================================================================
s = new_slide()
eyebrow_title(s, "PRÉSENTATION DU PROJET", "Le contexte et la chaîne de bout en bout", INDIGO)
card(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(2.0), "Le besoin métier",
     ["Une entreprise e-commerce reçoit des centaines d'avis clients.",
      "Les lire un par un est lent et coûteux.",
      "Objectif : détecter automatiquement si un avis est positif ou négatif, en français."],
     accent=INDIGO)
card(s, Inches(0.6), Inches(3.85), Inches(5.9), Inches(2.6), "Ce que j'ai construit",
     ["Pipeline : collecte → nettoyage → base de données.",
      "Deux APIs REST (données + modèle IA) sécurisées par clé API.",
      "Une application web Streamlit pour l'utilisateur final.",
      "Modèle XGBoost FR servi par l'API, suivi dans MLflow.",
      "Le tout testé, conteneurisé (Docker) et supervisé (Prometheus/Grafana)."],
     accent=VIOLET)
# Chaîne de traitement (5 étapes)
steps = ["Collecte (CSV · API REST · scraping)", "Nettoyage & RGPD (doublons, anonymisation)",
         "Base de données SQLite (Merise)", "APIs REST FastAPI (données + modèle IA)",
         "Application Streamlit (analyse temps réel)"]
y = Inches(1.7)
for i, st in enumerate(steps):
    c = [INDIGO, VIOLET, TEAL, ORANGE, CRIMSON][i]
    num = rect(s, Inches(6.9), y, Inches(0.5), Inches(0.5), c, shape=MSO_SHAPE.OVAL)
    tf = num.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, str(i+1), 14, WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, Inches(7.6), y + Inches(0.05), Inches(5.2), Inches(0.5),
            [(st, 13, TEXTDARK, True, PP_ALIGN.LEFT)])
    y += Inches(0.95)
footer(s, _page["n"])

# =================================================================
# 3. PROBLÉMATIQUE (3 stats)
# =================================================================
s = new_slide()
eyebrow_title(s, "PROBLÉMATIQUE", "Pourquoi automatiser l'analyse des avis ?", INDIGO)
stats = [("100s", "d'avis / mois", CRIMSON), ("2 classes", "Positif / Négatif", ORANGE),
         ("< 2 s", "par prédiction", INDIGO)]
x = Inches(0.7)
for big, small, c in stats:
    b = rect(s, x, Inches(1.8), Inches(3.85), Inches(1.7), WHITE, line=c, line_w=Pt(2),
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x, Inches(2.0), Inches(3.85), Inches(0.9),
            [(big, 36, c, True, PP_ALIGN.CENTER)])
    textbox(s, x, Inches(2.95), Inches(3.85), Inches(0.5),
            [(small, 14, GREY, False, PP_ALIGN.CENTER)])
    x += Inches(4.1)
card(s, Inches(0.7), Inches(3.9), Inches(11.85), Inches(2.3), "Le constat",
     ["Lire et trier manuellement des centaines d'avis est lent, coûteux et peu reproductible.",
      "Le besoin : un service fiable qui classe un avis en français en moins de 2 secondes.",
      "Enjeux : confidentialité des données (RGPD), souveraineté (modèle local) et traçabilité du modèle.",
      "Cible : e-commerçants, service client, équipes produit et analystes."],
     accent=INDIGO)
footer(s, _page["n"])

# =================================================================
# 4. PLAN (5 épreuves)
# =================================================================
s = new_slide()
eyebrow_title(s, "PLAN", "Les 5 épreuves de la certification", INDIGO)
rows = [("E1", "Collecte, stockage & exposition des données", "C1 → C5", "15 min", INDIGO),
        ("E3", "Le modèle IA et son exposition + démo", "C9 → C13", "20 min", VIOLET),
        ("E4", "Conception & industrialisation de l'app + démo", "C14 → C19", "20 min", TEAL),
        ("E2", "Veille, benchmark & cahier des charges", "C6 → C8", "15 min", ORANGE),
        ("E5", "Supervision & gestion des incidents", "C20 → C21", "10 min", CRIMSON)]
y = Inches(1.75)
for ep, title, comp, dur, c in rows:
    rect(s, Inches(0.6), y, Inches(12.13), Inches(0.85), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    blk = rect(s, Inches(0.7), y + Inches(0.13), Inches(1.3), Inches(0.6), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = blk.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, ep, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, Inches(2.2), y + Inches(0.22), Inches(7.5), Inches(0.5),
            [(title, 15, TEXTDARK, True, PP_ALIGN.LEFT)])
    textbox(s, Inches(9.8), y + Inches(0.25), Inches(1.6), Inches(0.4),
            [(comp, 12, GREY, False, PP_ALIGN.LEFT)])
    textbox(s, Inches(11.4), y + Inches(0.22), Inches(1.2), Inches(0.4),
            [(dur, 14, c, True, PP_ALIGN.RIGHT)])
    y += Inches(0.97)
footer(s, _page["n"])

# =================================================================
# E1
# =================================================================
section_slide("E1", "Collecte, stockage & exposition des données", "Compétences C1 à C5")

# E1 — objectifs
s = new_slide(); chip(s, "E1", INDIGO)
eyebrow_title(s, "E1 · OBJECTIFS", "Objectifs et compétences couvertes", INDIGO)
defs = [("C1", "Extraction multi-sources", "CSV + API REST + scraping BeautifulSoup"),
        ("C2", "Requêtes SQL", "Jointures, agrégations, requêtes paramétrées"),
        ("C3", "Nettoyage & agrégation", "Doublons, regex, RGPD, features qualité"),
        ("C4", "Base de données", "Modélisation Merise (MCD/MLD/MPD), SQLite, RGPD"),
        ("C5", "API REST Data", "FastAPI, endpoints documentés Swagger")]
xs = [Inches(0.6), Inches(4.65), Inches(8.7)]
for i, (cc, t, d) in enumerate(defs):
    col = i % 3; row = i // 3
    x = xs[col]; y = Inches(1.8) + row * Inches(2.3)
    card(s, x, y, Inches(3.85), Inches(2.05), t, [d], accent=INDIGO, badge=cc)
footer(s, _page["n"])

# E1 C1
s = new_slide(); chip(s, "E1", INDIGO)
eyebrow_title(s, "E1 · COLLECTE", "C1 — Extraction depuis des sources hétérogènes", INDIGO)
card(s, Inches(0.6), Inches(1.8), Inches(3.85), Inches(3.6), "3 sources collectées",
     ["Fichier CSV local (avis type Amazon / Fnac).",
      "API REST publique (requests).",
      "Scraping web (BeautifulSoup).",
      "Harmonisation vers un schéma unique.",
      "Fusion dans un seul DataFrame pandas."], accent=INDIGO, badge="C1")
card(s, Inches(4.65), Inches(1.8), Inches(3.85), Inches(3.6), "Bonnes pratiques",
     ["RGPD dès la collecte : données publiques et anonymes.",
      "Sauvegarde des données brutes avant traitement.",
      "Gestion des erreurs réseau (source ignorée si KO).",
      "Code documenté et fonctions réutilisables."], accent=GREEN, title_color=GREEN)
image_placeholder(s, Inches(8.7), Inches(1.8), Inches(4.0), Inches(3.6),
                  "Schéma des 3 sources / extrait de code collecte")
footer(s, _page["n"])

# E1 C2
s = new_slide(); chip(s, "E1", INDIGO)
eyebrow_title(s, "E1 · COLLECTE", "C2 — Requêtes SQL d'extraction", INDIGO)
card(s, Inches(0.6), Inches(1.8), Inches(5.5), Inches(2.0), "Lecture & jointures",
     ["SELECT + JOIN entre avis et produits.",
      "Filtres produit, recherche LIKE, pagination LIMIT/OFFSET."], accent=INDIGO, badge="C2")
card(s, Inches(0.6), Inches(3.95), Inches(5.5), Inches(2.1), "Agrégations",
     ["GROUP BY produit / source.",
      "COUNT, AVG, MIN, MAX, SUM.",
      "% d'avis positifs (CASE WHEN)."], accent=TEAL, title_color=TEAL)
code_box(s, Inches(6.4), Inches(1.8), Inches(6.3), Inches(4.25),
         "Requête paramétrée (anti-injection SQL)",
         "SELECT p.nom AS produit,\n"
         "       COUNT(a.id)       AS nb_avis,\n"
         "       ROUND(AVG(a.note),2) AS note_moy,\n"
         "       SUM(CASE WHEN a.note>=4\n"
         "           THEN 1 ELSE 0 END) AS positifs\n"
         "FROM produits p\n"
         "LEFT JOIN avis a ON a.produit_id = p.id\n"
         "GROUP BY p.id\n"
         "ORDER BY note_moy DESC;\n\n"
         "# paramètre préparé -> protège l'injection\n"
         "conn.execute(\n"
         '  "... WHERE p.nom = ?", (nom_produit,))')
footer(s, _page["n"])

# E1 C3
s = new_slide(); chip(s, "E1", INDIGO)
eyebrow_title(s, "E1 · COLLECTE", "C3 — Nettoyage & agrégation des données", INDIGO)
clean = [("1", "Doublons", "Suppression sur le texte (drop_duplicates)."),
         ("2", "Texte", "Regex, espaces, caractères spéciaux, max 1000 car."),
         ("3", "Notes", "Validation : forcées dans la plage 1 à 5."),
         ("4", "RGPD", "Anonymisation auteurs : jean.dupont → jea*****."),
         ("5", "Dates", "Normalisation au format ISO YYYY-MM-DD."),
         ("6", "Vides", "Suppression des avis < 10 caractères.")]
xs = [Inches(0.6), Inches(4.65), Inches(8.7)]
for i, (n, t, d) in enumerate(clean):
    col = i % 3; row = i // 3
    x = xs[col]; y = Inches(1.8) + row * Inches(1.85)
    box = rect(s, x, y, Inches(3.85), Inches(1.6), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    num = rect(s, x + Inches(0.2), y + Inches(0.2), Inches(0.45), Inches(0.45),
               [INDIGO, VIOLET, TEAL, ORANGE, CRIMSON, RGBColor(0x0E,0xA5,0xE9)][i], shape=MSO_SHAPE.OVAL)
    tf = num.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set_text(tf, n, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.8), y + Inches(0.18), Inches(2.9), Inches(0.4),
            [(t, 14, TEXTDARK, True, PP_ALIGN.LEFT)])
    textbox(s, x + Inches(0.25), y + Inches(0.7), Inches(3.4), Inches(0.8),
            [(d, 11, GREY, False, PP_ALIGN.LEFT)])
banner(s, "Puis une agrégation par produit (nb d'avis, note moyenne, % positifs) = 1er contrôle qualité", INDIGO)
footer(s, _page["n"])

# E1 C4
s = new_slide(); chip(s, "E1", INDIGO)
eyebrow_title(s, "E1 · STOCKAGE BDD", "C4 — Modélisation Merise & base de données", INDIGO)
card(s, Inches(0.6), Inches(1.8), Inches(3.85), Inches(3.5), "Démarche Merise",
     ["MCD — entités métier et cardinalités.",
      "MLD — tables relationnelles, clés primaires/étrangères.",
      "MPD — implémentation SQLite (schema.sql)."], accent=INDIGO, badge="C4")
card(s, Inches(4.65), Inches(1.8), Inches(3.85), Inches(3.5), "La base (SQLite)",
     ["5 tables : produits, avis, predictions_sentiment, logs_api, consentements_rgpd.",
      "Contraintes : CHECK note 1-5, clés étrangères.",
      "Index + vue avis_avec_sentiment."], accent=TEAL, title_color=TEAL)
card(s, Inches(8.7), Inches(1.8), Inches(4.0), Inches(3.5), "Conformité RGPD",
     ["Minimisation : ni email, ni IP, ni nom complet.",
      "Pseudonymisation des auteurs.",
      "Durée de conservation (a_supprimer_le) + droit à l'oubli.",
      "Table de consentements (Art. 30)."], accent=GREEN, title_color=GREEN)
image_placeholder(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.55),
                  "Diagramme MCD / MLD (depuis l'artefact)")
footer(s, _page["n"])

# E1 C5
s = new_slide(); chip(s, "E1", INDIGO)
eyebrow_title(s, "E1 · API REST DATA", "C5 — API REST des données avec FastAPI", INDIGO)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.0), "Ce que fait l'API",
     ["Expose avis et statistiques (HTTP/JSON).",
      "Endpoints : /avis, /statistiques, /recherche, /produits, /logs.",
      "Documentation Swagger automatique (/docs)."], accent=INDIGO, badge="C5")
card(s, Inches(0.6), Inches(3.95), Inches(5.6), Inches(2.1), "Sécurité & robustesse",
     ["Authentification par clé API (X-API-Key).",
      "Validation Pydantic des entrées/sorties.",
      "Middleware qui journalise chaque requête."], accent=CRIMSON, title_color=CRIMSON)
code_box(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(2.6), "Un endpoint protégé",
         '@app.get("/avis")\n'
         "def lister_avis(\n"
         "  limit: int = Query(20, ge=1, le=100),\n"
         "  offset: int = 0,\n"
         "  _key: str = Depends(verifier_api_key)):\n"
         "  return queries.get_tous_les_avis(...)")
image_placeholder(s, Inches(6.5), Inches(4.55), Inches(6.2), Inches(1.5),
                  "Capture Swagger /docs (API Données)")
footer(s, _page["n"])

# =================================================================
# E3 — Modèle IA
# =================================================================
section_slide("E3", "Le modèle IA et son exposition", "Compétences C9 à C13  ·  avec démo")

# E3 objectifs
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · OBJECTIFS", "Objectifs et compétences couvertes", VIOLET)
defs = [("C9", "API du modèle", "Endpoint /predict exposant le modèle (FastAPI)"),
        ("C10", "Intégration", "Branchement de l'API dans l'app Streamlit"),
        ("C11", "Modèle & monitoring", "Entraînement, métriques, MLflow, suivi"),
        ("C12", "Tests automatisés", "Suite pytest : modèle, APIs, logique métier"),
        ("C13", "Livraison continue", "Git + CI/CD avec rapport d'évaluation")]
xs = [Inches(0.6), Inches(4.65), Inches(8.7)]
for i, (cc, t, d) in enumerate(defs):
    col = i % 3; row = i // 3
    card(s, xs[col], Inches(1.8) + row*Inches(2.3), Inches(3.85), Inches(2.05),
         t, [d], accent=VIOLET, badge=cc)
footer(s, _page["n"])

# E3 — choix modèle (LE point clé : XGBoost FR cohérent)
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · LE MODÈLE CHOISI", "C11 — Du DistilBERT anglais au XGBoost français", VIOLET)
card(s, Inches(0.6), Inches(1.8), Inches(5.9), Inches(2.1), "Le problème de cohérence",
     ["L'API servait DistilBERT (anglais, SST-2)…",
      "…alors que le modèle entraîné était XGBoost en français.",
      "Incohérence : prédictions peu fiables sur des avis FR."],
     accent=CRIMSON, title_color=CRIMSON, badge="C11")
card(s, Inches(0.6), Inches(4.0), Inches(5.9), Inches(2.2), "La solution retenue",
     ["Modèle de production unifié : XGBoost (TF-IDF) entraîné en FR.",
      "Servi par l'API /predict ET utilisé par le front.",
      "Léger (~1,5 Mo), CPU pur, souverain (100 % local, RGPD).",
      "Accuracy 0,84 · F1 0,84 · ROC-AUC 0,93 (jeu de test FR)."],
     accent=GREEN, title_color=GREEN)
image_placeholder(s, Inches(6.7), Inches(1.8), Inches(6.0), Inches(4.4),
                  "Matrice de confusion XGBoost FR (depuis MLflow)")
footer(s, _page["n"])

# E3 — benchmark des modèles FR (notre benchmark)
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · BENCHMARK MODÈLES", "C11 — Benchmark des modèles de sentiment FR", VIOLET)
card(s, Inches(0.6), Inches(1.75), Inches(5.4), Inches(4.45), "Protocole comparatif",
     ["Même jeu de test français (split seed 42).",
      "XGBoost FR vs modèles transformers, en zero-shot.",
      "Métriques : accuracy, F1, ROC-AUC + latence.",
      "Chaque modèle = 1 run MLflow ('benchmark-sentiment-fr').",
      "",
      "3 enseignements :",
      "1) L'ancien DistilBERT anglais s'effondre en FR (F1 0,38).",
      "2) nlptown (multilingue) un peu plus précis… mais ~600× plus lent.",
      "3) XGBoost FR = meilleur compromis précision / latence / RGPD."],
     accent=VIOLET, badge="C11")
# Tableau des résultats réels (n_test=300)
bcols = ["Modèle", "F1", "ROC-AUC", "ms/avis"]
brows = [("nlptown (multilingue)", "0.857", "0.949", "794", GREY),
         ("XGBoost (FR) — retenu", "0.831", "0.917", "1.3", GREEN),
         ("distilbert-en (ancien)", "0.381", "0.649", "679", CRIMSON)]
bx = [Inches(6.2), Inches(9.1), Inches(10.3), Inches(11.6)]
bw = [Inches(2.9), Inches(1.2), Inches(1.3), Inches(1.1)]
rect(s, Inches(6.2), Inches(1.75), Inches(6.5), Inches(0.55), NAVY)
for i, c in enumerate(bcols):
    textbox(s, bx[i] + Inches(0.1), Inches(1.83), bw[i], Inches(0.4),
            [(c, 12, WHITE, True, PP_ALIGN.LEFT)])
yy = Inches(2.32)
for name, f1, roc, ms, col in brows:
    rect(s, Inches(6.2), yy, Inches(6.5), Inches(0.7), LIGHT if col != GREEN else RGBColor(0xE6,0xF7,0xEE))
    textbox(s, bx[0] + Inches(0.1), yy + Inches(0.16), bw[0], Inches(0.4),
            [(name, 12, col, True, PP_ALIGN.LEFT)])
    for j, v in enumerate([f1, roc, ms]):
        textbox(s, bx[j+1] + Inches(0.1), yy + Inches(0.16), bw[j+1], Inches(0.4),
                [(v, 12, TEXTDARK, name.startswith("XGBoost"), PP_ALIGN.LEFT)])
    yy += Inches(0.72)
image_placeholder(s, Inches(6.2), Inches(4.5), Inches(6.5), Inches(1.7),
                  "Graphe F1 par modèle (benchmark_f1.png, MLflow)")
banner(s, "XGBoost FR : F1 0,83 · ROC-AUC 0,92 · 1,3 ms/avis — ~600× plus rapide que les transformers, 100% local", VIOLET)
footer(s, _page["n"])

# E3 — MLflow
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · SUIVI", "C11 — Fine-tuning (GridSearch), learning curve & MLflow", VIOLET)
card(s, Inches(0.6), Inches(1.8), Inches(5.5), Inches(2.0), "Optimisation",
     ["GridSearchCV + validation croisée.",
      "Grille : n_estimators, max_depth, learning_rate.",
      "On garde best_params_ et best_score_."], accent=VIOLET, badge="C11")
card(s, Inches(0.6), Inches(3.95), Inches(5.5), Inches(2.1), "Suivi MLflow (dockerisé)",
     ["Serveur MLflow dans docker-compose (port 5000).",
      "Chaque run loggue params, métriques et artefacts.",
      "Modèle + courbes + rapport sauvegardés et rechargeables."], accent=TEAL, title_color=TEAL)
image_placeholder(s, Inches(6.4), Inches(1.8), Inches(6.3), Inches(4.25),
                  "Capture UI MLflow (runs / artefacts / learning curve)")
banner(s, "DÉMO — Ouvrir l'UI MLflow : comparer les runs, métriques et hyperparamètres", VIOLET)
footer(s, _page["n"])

# E3 C9 API predict
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · API REST IA", "C9 — API REST exposant le modèle IA", VIOLET)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.0), "Endpoints exposés",
     ["POST /predict — sentiment d'un texte (label + score + durée).",
      "POST /predict/batch — jusqu'à 10 textes.",
      "GET /monitoring — statistiques d'utilisation."], accent=VIOLET, badge="C9")
card(s, Inches(0.6), Inches(3.95), Inches(5.6), Inches(2.1), "Sécurité — OWASP API Top 10",
     ["Clé API obligatoire sur les routes sensibles (A01/A02).",
      "Validation Pydantic des entrées (A03).",
      "Limites sur le batch pour éviter les abus (A04)."], accent=CRIMSON, title_color=CRIMSON)
code_box(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.25), "Réponse de /predict",
         'POST /predict\n'
         '{ "texte": "Ce produit est excellent !" }\n\n'
         "->{\n"
         '  "label": "POSITIVE",\n'
         '  "score": 0.97,\n'
         '  "sentiment_fr": "Positif",\n'
         '  "duree_ms": 12.6,\n'
         '  "modele": "XGBoost (TF-IDF) - FR"\n'
         "}")
footer(s, _page["n"])

# E3 DÉMO
s = prs.slides.add_slide(BLANK); _page["n"] += 1
rect(s, 0, 0, SW, SH, NAVY)
oval = rect(s, Inches(6.07), Inches(1.6), Inches(1.2), Inches(1.2), VIOLET, shape=MSO_SHAPE.OVAL)
tf = oval.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_set_text(tf, "▶", 36, NAVY, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(0), Inches(3.0), Inches(13.33), Inches(0.8),
        [("DÉMONSTRATION LIVE — Modèle & API", 30, WHITE, True, PP_ALIGN.CENTER)])
textbox(s, Inches(0), Inches(4.0), Inches(13.33), Inches(1.2),
        [("1. Swagger /docs → s'authentifier (X-API-Key)   2. POST /predict sur un avis FR positif puis négatif",
          15, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.CENTER),
         ("3. Réponse 200 (label + score)   4. Entrée invalide → 422 (validation Pydantic)",
          15, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.CENTER)])
footer(s, _page["n"])

# E3 C10 integration
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · INTÉGRATION", "C10 — Intégrer le modèle dans l'application", VIOLET)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(4.25), "Du modèle à l'utilisateur",
     ["L'app Streamlit n'embarque pas le modèle : elle appelle /predict.",
      "Architecture découplée : changer de modèle sans toucher au front.",
      "La clé API est envoyée à chaque appel (champ dans la sidebar).",
      "Affichage : sentiment, score de confiance, temps de traitement.",
      "Mode batch : analyser plusieurs avis d'un coup."], accent=VIOLET, badge="C10")
code_box(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.25), "Appel de l'API depuis Streamlit",
         "resultat = appeler_api(\n"
         '  f"{API_MODELE_URL}/predict",\n'
         '  methode="POST",\n'
         '  donnees={"texte": texte})\n\n'
         'if resultat["label"] == "POSITIVE":\n'
         '  st.success(f"Positif - "\n'
         "    f\"{resultat['score']*100:.1f}%\")\n"
         "else:\n"
         '  st.error("Négatif")')
footer(s, _page["n"])

# E3 C12 tests
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · TESTS", "C12 — Tests automatisés avec pytest", VIOLET)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.1), "Ce que je teste",
     ["Les 2 APIs (test_model_api, test_data_api).",
      "Le modèle XGBoost : pipeline, métriques, GridSearch, learning curve."], accent=VIOLET, badge="C12")
card(s, Inches(0.6), Inches(4.05), Inches(5.6), Inches(2.0), "Bonnes pratiques",
     ["Fixtures : jeu de données jouet déterministe.",
      "Tests rapides, isolés et reproductibles.",
      "Couverture de code mesurée (pytest-cov)."], accent=GREEN, title_color=GREEN)
code_box(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.25), "Exécution",
         "$ pytest tests/ -v --cov=.\n\n"
         "test_data_api.py ......   [ 33%]\n"
         "test_model_api.py .....   [ 66%]\n"
         "test_xgboost_sentiment.py [100%]\n\n"
         "=== passed in 6.4s ===\n\n"
         "def test_metrics_keys():\n"
         "  m = xs.evaluate_classification(...)\n"
         '  assert "f1" in m and "roc_auc" in m')
footer(s, _page["n"])

# E3 C13 CI
s = new_slide(); chip(s, "E3", VIOLET)
eyebrow_title(s, "E3 · VERSIONNAGE", "C13 — Push GitHub & intégration continue", VIOLET)
card(s, Inches(0.6), Inches(1.8), Inches(3.85), Inches(3.5), "Git & GitHub",
     ["Code versionné sur GitHub.",
      "Commits conventionnels (feat:, fix:, test:).",
      "Branches main / develop + Pull Requests.",
      ".gitignore (cache, secrets, données lourdes)."], accent=VIOLET, badge="C13")
card(s, Inches(4.65), Inches(1.8), Inches(3.85), Inches(3.5), "CI au push",
     ["Workflow GitHub Actions à chaque push / PR.",
      "Installe les dépendances, initialise la BDD.",
      "Valide données + modèle, lance pytest."], accent=ORANGE, title_color=ORANGE)
card(s, Inches(8.7), Inches(1.8), Inches(4.0), Inches(3.5), "Pourquoi c'est important",
     ["Historique complet et traçable.",
      "Collaboration et revue de code.",
      "Chaque modification testée avant fusion."], accent=GREEN, title_color=GREEN)
image_placeholder(s, Inches(0.6), Inches(5.45), Inches(12.1), Inches(1.55),
                  "Capture GitHub Actions (workflow vert)")
footer(s, _page["n"])

# =================================================================
# E4 — Application
# =================================================================
section_slide("E4", "Conception & industrialisation de l'application", "Compétences C14 à C19  ·  avec démo")

# E4 objectifs
s = new_slide(); chip(s, "E4", TEAL)
eyebrow_title(s, "E4 · OBJECTIFS", "Objectifs et compétences couvertes", TEAL)
defs = [("C14", "Cadre technique", "Stack et choix technologiques"),
        ("C15", "Architecture", "Front / API / base de données"),
        ("C16-17", "Développer", "Agile + interface Streamlit"),
        ("C18", "Tester", "Tests applicatifs et CI"),
        ("C19", "Livrer", "Docker + déploiement")]
xs = [Inches(0.6), Inches(4.65), Inches(8.7)]
for i, (cc, t, d) in enumerate(defs):
    col = i % 3; row = i // 3
    card(s, xs[col], Inches(1.8) + row*Inches(2.3), Inches(3.85), Inches(2.05),
         t, [d], accent=TEAL, title_color=TEAL, badge=cc)
footer(s, _page["n"])

# E4 C14 stack
s = new_slide(); chip(s, "E4", TEAL)
eyebrow_title(s, "E4 · CADRAGE", "C14 — Les technologies utilisées", TEAL)
cats = [("Développement", "VS Code · Jupyter · GitHub · Trello"),
        ("Données", "CSV · API REST · BeautifulSoup"),
        ("Modélisation", "scikit-learn · XGBoost · MLflow"),
        ("Base de données", "SQLite · SQL"),
        ("API & Interface", "FastAPI · Pydantic · Streamlit"),
        ("Déploiement & Monitoring", "Docker · Prometheus · Grafana · pytest")]
x = Inches(0.6); w = Inches(1.95)
for t, d in cats:
    rect(s, x, Inches(1.9), w, Inches(4.2), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x + Inches(0.1), Inches(2.05), w - Inches(0.2), Inches(0.9),
            [(t, 12, TEAL, True, PP_ALIGN.CENTER)])
    image_placeholder(s, x + Inches(0.15), Inches(3.0), w - Inches(0.3), Inches(1.4), "logos")
    textbox(s, x + Inches(0.1), Inches(4.5), w - Inches(0.2), Inches(1.5),
            [(d, 10, GREY, False, PP_ALIGN.CENTER)])
    x += Inches(2.02)
footer(s, _page["n"])

# E4 C15 architecture
s = new_slide(); chip(s, "E4", TEAL)
eyebrow_title(s, "E4 · CADRAGE", "C15 — Architecture : 6 services conteneurisés", TEAL)
svcs = [("Streamlit", "8501", "Frontend", TEAL), ("API Modèle", "8000", "XGBoost FR", ORANGE),
        ("API Données", "8001", "BDD/SQL", INDIGO), ("Prometheus", "9090", "Métriques", CRIMSON),
        ("Grafana", "3000", "Dashboards", VIOLET), ("MLflow", "5000", "Tracking", GREEN)]
x = Inches(0.6); y = Inches(2.0)
for i, (n, port, role, c) in enumerate(svcs):
    col = i % 3; row = i // 3
    xx = Inches(0.6) + col * Inches(4.1); yy = Inches(2.0) + row * Inches(2.0)
    box = rect(s, xx, yy, Inches(3.85), Inches(1.7), WHITE, line=c, line_w=Pt(2),
               shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, xx, yy, Inches(3.85), Inches(0.55), c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = textbox(s, xx, yy + Inches(0.06), Inches(3.85), Inches(0.45),
                 [(n, 15, WHITE, True, PP_ALIGN.CENTER)])
    textbox(s, xx, yy + Inches(0.7), Inches(3.85), Inches(0.4),
            [("Port " + port, 13, TEXTDARK, True, PP_ALIGN.CENTER)])
    textbox(s, xx, yy + Inches(1.15), Inches(3.85), Inches(0.4),
            [(role, 12, GREY, False, PP_ALIGN.CENTER)])
banner(s, "Architecture découplée : chaque service est isolé, remplaçable et orchestré par Docker Compose", TEAL)
footer(s, _page["n"])

# E4 C16 agile
s = new_slide(); chip(s, "E4", TEAL)
eyebrow_title(s, "E4 · SUIVI AGILE", "C16 — Coordination en méthode Agile", TEAL)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.5), "4 sprints",
     ["Sprint 1 — Données (collecte, nettoyage, BDD).",
      "Sprint 2 — APIs (données + modèle IA).",
      "Sprint 3 — Frontend + tests + CI.",
      "Sprint 4 — Docker, monitoring, documentation."], accent=TEAL, title_color=TEAL, badge="C16")
card(s, Inches(0.6), Inches(4.45), Inches(5.6), Inches(1.7), "Rituels & outils",
     ["GitHub Projects (Kanban) + Issues + Pull Requests.",
      "Daily, sprint review, rétrospective."], accent=ORANGE, title_color=ORANGE)
image_placeholder(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.35),
                  "Capture board Kanban (Trello / GitHub Projects)")
footer(s, _page["n"])

# E4 C17 app
s = new_slide(); chip(s, "E4", TEAL)
eyebrow_title(s, "E4 · APPLICATION", "C17 — Application Streamlit (en français)", TEAL)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.4), "Les pages de l'app",
     ["Analyser un texte — saisie libre + résultat instantané.",
      "Tableau de bord — métriques globales et graphiques.",
      "Base de données — exploration, filtres, recherche, batch.",
      "Authentification par clé API dans la barre latérale."], accent=TEAL, title_color=TEAL, badge="C17")
card(s, Inches(0.6), Inches(4.35), Inches(5.6), Inches(1.8), "Conception & qualité",
     ["Wireframes et parcours définis avant le code.",
      "Accessibilité WCAG 2.1 AA (contrastes, labels).",
      "Interface 100 % en français, modèle XGBoost FR."], accent=GREEN, title_color=GREEN)
image_placeholder(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.35),
                  "Capture de l'app Streamlit (page Analyser)")
footer(s, _page["n"])

# E4 DÉMO
s = prs.slides.add_slide(BLANK); _page["n"] += 1
rect(s, 0, 0, SW, SH, NAVY)
oval = rect(s, Inches(6.07), Inches(1.8), Inches(1.2), Inches(1.2), TEAL, shape=MSO_SHAPE.OVAL)
tf = oval.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
_set_text(tf, "▶", 36, NAVY, bold=True, align=PP_ALIGN.CENTER)
textbox(s, Inches(0), Inches(3.2), Inches(13.33), Inches(0.8),
        [("DÉMONSTRATION LIVE — Application", 30, WHITE, True, PP_ALIGN.CENTER)])
textbox(s, Inches(0), Inches(4.2), Inches(13.33), Inches(0.8),
        [("Lancer l'app Streamlit → analyser un avis FR, parcourir le tableau de bord et la base.",
          15, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.CENTER)])
footer(s, _page["n"])

# E4 C18 CI / C19 docker
s = new_slide(); chip(s, "E4", TEAL)
eyebrow_title(s, "E4 · CI/CD", "C18-C19 — Tests automatisés, Docker & livraison", TEAL)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(4.25), "Pipeline CI + conteneurisation",
     ["GitHub Actions à chaque push / PR (lint, BDD, pytest, build).",
      "Un Dockerfile par service + docker-compose.",
      "Image API modèle allégée (XGBoost, sans torch).",
      "requirements séparés : complet / API modèle / API données.",
      "Versions figées pour la reproductibilité."], accent=TEAL, title_color=TEAL, badge="C18-19")
code_box(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.25), "Lancer toute la stack",
         "$ docker compose up --build\n\n"
         "OK api_modele   :8000  (XGBoost FR)\n"
         "OK api_donnees  :8001\n"
         "OK app (Streamlit):8501\n"
         "OK prometheus   :9090\n"
         "OK grafana      :3000\n"
         "OK mlflow       :5000")
footer(s, _page["n"])

# =================================================================
# E2 — Veille
# =================================================================
section_slide("E2", "Veille, benchmark & cahier des charges", "Compétences C6 à C8")

# E2 objectifs
s = new_slide(); chip(s, "E2", ORANGE)
eyebrow_title(s, "E2 · OBJECTIFS", "Objectifs et compétences couvertes", ORANGE)
defs = [("C6", "Veille technique", "Veille planifiée + analyse budgétaire, sources recoupées"),
        ("C7", "Identifier le service", "Benchmark services existants vs modèle sur mesure"),
        ("C8", "Paramétrer le service", "Environnement Docker, dépendances, POC documenté")]
x = Inches(0.6)
for cc, t, d in defs:
    card(s, x, Inches(2.0), Inches(3.95), Inches(3.0), t, [d], accent=ORANGE, title_color=ORANGE, badge=cc)
    x += Inches(4.05)
footer(s, _page["n"])

# E2 C6 veille
s = new_slide(); chip(s, "E2", ORANGE)
eyebrow_title(s, "E2 · VEILLE", "C6 — Veille technologique & analyse budgétaire", ORANGE)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.0), "Organisation de la veille",
     ["Créneaux hebdomadaires (~1h/semaine).",
      "Feedly (RSS), newsletters, GitHub Watch, LinkedIn.",
      "Sources : HuggingFace, arXiv, Papers With Code, CNIL, AI Act."], accent=ORANGE, title_color=ORANGE, badge="C6")
card(s, Inches(0.6), Inches(3.95), Inches(5.6), Inches(2.1), "Veille réglementaire",
     ["RGPD : minimisation, droit à l'oubli, conservation.",
      "AI Act : système à risque minimal (transparence)."], accent=GREEN, title_color=GREEN)
# table budget
hdr = ["Solution", "Coût"]
data = [("Modèle local (XGBoost / HuggingFace)", "0 € — open source"),
        ("AWS Comprehend", "~0,01 $ / 1000 unités"),
        ("Google NLP", "~0,001 $ / requête"),
        ("Azure Text Analytics", "~0,001 $ / requête"),
        ("OpenAI GPT-4", "élevé + hors UE")]
tb = rect(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(0.5), NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, Inches(6.7), Inches(1.88), Inches(4), Inches(0.4), [("Solution", 12, WHITE, True, PP_ALIGN.LEFT)])
textbox(s, Inches(10.6), Inches(1.88), Inches(2), Inches(0.4), [("Coût", 12, WHITE, True, PP_ALIGN.LEFT)])
y = Inches(2.35)
for i, (sol, cout) in enumerate(data):
    fill = LIGHT if i % 2 == 0 else WHITE
    rect(s, Inches(6.5), y, Inches(6.2), Inches(0.55), fill)
    textbox(s, Inches(6.7), y + Inches(0.1), Inches(4.0), Inches(0.4),
            [(sol, 11, TEXTDARK, sol.startswith("Modèle"), PP_ALIGN.LEFT)])
    textbox(s, Inches(10.6), y + Inches(0.1), Inches(2.0), Inches(0.4),
            [(cout, 11, ORANGE if i==0 else GREY, i==0, PP_ALIGN.LEFT)])
    y += Inches(0.57)
banner(s, "Choix : modèle local — coût nul + souveraineté/RGPD totale (les avis ne quittent pas le serveur)", ORANGE)
footer(s, _page["n"])

# E2 C7 benchmark services
s = new_slide(); chip(s, "E2", ORANGE)
eyebrow_title(s, "E2 · BENCHMARK", "C7 — Benchmark des services IA", ORANGE)
cols = ["Critère", "Local (XGBoost FR)", "AWS / Google / Azure", "OpenAI"]
rowsd = [("Prix", "Gratuit", "~0,001-0,01 $/req", "Élevé"),
         ("Précision (FR)", "~0,84 F1", "0,85-0,92", "Élevée"),
         ("Confidentialité", "Totale (local)", "Données → cloud", "Données → US"),
         ("RGPD", "Optimal", "Hors UE possible", "Hors UE"),
         ("Éco-impact", "Très faible (CPU)", "Moyen", "Élevé")]
# header
xcols = [Inches(0.6), Inches(3.4), Inches(6.6), Inches(10.0)]
wcols = [Inches(2.8), Inches(3.2), Inches(3.4), Inches(2.7)]
hb = rect(s, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.55), NAVY)
hl = rect(s, xcols[1], Inches(1.8), wcols[1], Inches(0.55), ORANGE)
for i, c in enumerate(cols):
    textbox(s, xcols[i] + Inches(0.1), Inches(1.88), wcols[i], Inches(0.4),
            [(c, 12, WHITE, True, PP_ALIGN.LEFT)])
y = Inches(2.4)
for r in rowsd:
    rect(s, Inches(0.6), y, Inches(12.1), Inches(0.7), LIGHT)
    for i, val in enumerate(r):
        textbox(s, xcols[i] + Inches(0.1), y + Inches(0.18), wcols[i], Inches(0.4),
                [(val, 12, TEXTDARK, i==1, PP_ALIGN.LEFT)])
    y += Inches(0.72)
banner(s, "Décision : modèle local XGBoost FR — souverain, gratuit, précis et sobre pour une classification binaire", ORANGE)
footer(s, _page["n"])

# E2 C8 POC
s = new_slide(); chip(s, "E2", ORANGE)
eyebrow_title(s, "E2 · CAHIER DES CHARGES", "C8 — Paramétrer & cadrer le service IA", ORANGE)
card(s, Inches(0.6), Inches(1.8), Inches(3.85), Inches(3.5), "Paramétrage",
     ["Pipeline TF-IDF + XGBoost (joblib).",
      "device CPU, modèle léger (~1,5 Mo).",
      "Chargement au démarrage de l'API.",
      "Variables d'env (clé API, URI MLflow)."], accent=ORANGE, title_color=ORANGE, badge="C8")
card(s, Inches(4.65), Inches(1.8), Inches(3.85), Inches(3.5), "Exigences",
     ["Fonctionnelles : analyse < 2s, stats, recherche, batch.",
      "Non-fonctionnelles : RGPD, dispo > 95%, Swagger, API Key, WCAG, OWASP."], accent=INDIGO)
card(s, Inches(8.7), Inches(1.8), Inches(4.0), Inches(3.5), "Preuve de concept",
     ["F1 ~0,84 (seuil > 0,80).",
      "Latence faible (CPU).",
      "Conformité RGPD totale.",
      "Coût 0 €, stack Docker fonctionnelle.",
      "→ Décision : passer en production."], accent=GREEN, title_color=GREEN)
footer(s, _page["n"])

# =================================================================
# E5 — Monitoring
# =================================================================
section_slide("E5", "Supervision & gestion des incidents", "Compétences C20 à C21")

# E5 objectifs
s = new_slide(); chip(s, "E5", CRIMSON)
eyebrow_title(s, "E5 · OBJECTIFS", "Objectifs et compétences couvertes", CRIMSON)
card(s, Inches(0.9), Inches(2.0), Inches(5.5), Inches(2.6), "C20 — Surveiller",
     ["Monitorage + journalisation (RGPD).",
      "Prometheus scrape /metrics toutes les 15s.",
      "Grafana : dashboards temps réel.",
      "Alerte sur seuils (erreurs, latence, drift)."], accent=CRIMSON, title_color=CRIMSON, badge="C20")
card(s, Inches(6.9), Inches(2.0), Inches(5.5), Inches(2.6), "C21 — Résoudre",
     ["Identifier la cause racine.",
      "Reproduire le problème.",
      "Corriger le code et la configuration.",
      "Documenter la solution (rapport d'incident)."], accent=ORANGE, title_color=ORANGE, badge="C21")
footer(s, _page["n"])

# E5 C20 monitoring
s = new_slide(); chip(s, "E5", CRIMSON)
eyebrow_title(s, "E5 · MONITORING", "C20 — Surveillance avec Prometheus & Grafana", CRIMSON)
card(s, Inches(0.6), Inches(1.8), Inches(5.6), Inches(2.2), "La chaîne de supervision",
     ["Chaque API expose /metrics (format Prometheus).",
      "Prometheus scrape toutes les 15s.",
      "Grafana : dashboards provisionnés automatiquement.",
      "MLflow (port 5000) pour le suivi des modèles."], accent=CRIMSON, title_color=CRIMSON, badge="C20")
card(s, Inches(0.6), Inches(4.15), Inches(5.6), Inches(1.95), "Métriques & seuils",
     ["Débit req/s, latence P50/P95, erreurs 4xx/5xx.",
      "Volume de prédictions et répartition des classes.",
      "Seuils WARNING / CRITIQUE documentés."], accent=ORANGE, title_color=ORANGE)
image_placeholder(s, Inches(6.5), Inches(1.8), Inches(6.2), Inches(4.3),
                  "Capture dashboard Grafana (rempli après trafic)")
footer(s, _page["n"])

# E5 C21 incident (RÉEL)
s = new_slide(); chip(s, "E5", CRIMSON)
eyebrow_title(s, "E5 · INCIDENT", "C21 — Incident résolu : « Grafana n'affiche rien »", CRIMSON)
card(s, Inches(0.6), Inches(1.7), Inches(5.9), Inches(2.3), "1. Détection & diagnostic",
     ["Dashboard Grafana vide sur tous les panels.",
      "Datasource OK, dashboard provisionné, APIs healthy.",
      "Requêtes PromQL renvoyaient des taux à 0 / NaN."], accent=CRIMSON, title_color=CRIMSON, badge="C21")
card(s, Inches(0.6), Inches(4.1), Inches(5.9), Inches(2.5), "3. Causes racines (3)",
     ["Aucun trafic : compteurs ~0 → rate() = 0 (rien à tracer).",
      "Conteneur Prometheus périmé : port 9090 non publié.",
      "Serveur MLflow bindé sur 127.0.0.1 (YAML mal plié) → injoignable depuis l'hôte."], accent=INDIGO)
card(s, Inches(6.7), Inches(1.7), Inches(6.0), Inches(2.3), "2. Reproduction",
     ["Vérification des cibles Prometheus (up).",
      "Inspection docker ps (ports) + docker logs.",
      "Test /health depuis l'hôte vs dans le conteneur."], accent=VIOLET, title_color=VIOLET)
card(s, Inches(6.7), Inches(4.1), Inches(6.0), Inches(2.5), "4. Correction + doc",
     ["Génération de trafic (generer_trafic.py) → dashboards remplis.",
      "Recréation du conteneur Prometheus (port publié).",
      "Commande MLflow sur une seule ligne (host 0.0.0.0).",
      "Image MLflow dédiée pour un redémarrage rapide."], accent=GREEN, title_color=GREEN)
footer(s, _page["n"])

# =================================================================
# CONCLUSION
# =================================================================
s = prs.slides.add_slide(BLANK); _page["n"] += 1
rect(s, 0, 0, SW, SH, NAVY)
textbox(s, Inches(0.8), Inches(0.7), Inches(11.5), Inches(1.0),
        [("Un projet IA complet, de la donnée à la production", 30, WHITE, True, PP_ALIGN.LEFT)])
kpis = [("21", "compétences couvertes (C1 → C21)", INDIGO),
        ("XGBoost FR", "modèle souverain, F1 ~0,84", VIOLET),
        ("6", "services orchestrés avec Docker", TEAL)]
x = Inches(0.8)
for big, small, c in kpis:
    box = rect(s, x, Inches(2.2), Inches(3.7), Inches(1.9), RGBColor(0x2A,0x27,0x5E),
               line=c, line_w=Pt(2), shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, x, Inches(2.45), Inches(3.7), Inches(0.9), [(big, 30, c, True, PP_ALIGN.CENTER)])
    textbox(s, x, Inches(3.35), Inches(3.7), Inches(0.7), [(small, 13, WHITE, False, PP_ALIGN.CENTER)])
    x += Inches(3.95)
textbox(s, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.6),
        [("Pistes d'évolution", 16, ORANGE, True, PP_ALIGN.LEFT)])
textbox(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.6),
        [("CamemBERT / fine-tuning métier · calibration des probabilités · alertes Grafana · déploiement cloud",
          14, RGBColor(0xC7,0xCB,0xF5), False, PP_ALIGN.LEFT)])
textbox(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.6),
        [("Merci de votre attention — je reste à votre disposition pour vos questions.",
          16, WHITE, True, PP_ALIGN.LEFT)])
footer(s, _page["n"])

# Merci
s = prs.slides.add_slide(BLANK); _page["n"] += 1
rect(s, 0, 0, SW, SH, NAVY)
textbox(s, Inches(0), Inches(3.1), Inches(13.33), Inches(1.2),
        [("Merci de votre écoute", 40, WHITE, True, PP_ALIGN.CENTER)])
image_placeholder(s, Inches(4.66), Inches(4.6), Inches(4.0), Inches(2.0), "Visuel de clôture (optionnel)")

out = "Presentation_Avis_Sentiment_AMELIOREE.pptx"
prs.save(out)
print(f"OK -> {out} | {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
